import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox
import threading
import os
import sys
import json
import re
import uuid
import time
from datetime import datetime
import traceback
import platform

# --- åŸºç¡€åº“ ---
try:
    import pyperclip
except ImportError:
    pyperclip = None

from openai import OpenAI
from duckduckgo_search import DDGS
import pypdf
from docx import Document

# --- é‡å‹åº“å®‰å…¨å¯¼å…¥ ---
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# --- é…ç½®åŒºåŸŸ ---
APP_NAME = "DeepSeek Pro"
APP_VERSION = "v2.5.2 (Fix Empty Bubble)"
DEV_NAME = "Yu Jinquan"

DEFAULT_CONFIG = {
    "api_key": "",
    "model": "deepseek-chat",
    "use_search": False,
    "is_r1": False,
    "system_prompt": "ä½ æ˜¯ä¸€ä¸ªä¹äºåŠ©äººçš„AIåŠ©æ‰‹ã€‚ä»£ç è¯·ç”¨Markdownæ ¼å¼ã€‚"
}

# é¢œè‰²é…ç½®
COLOR_USER_BUBBLE = "#95EC69" 
COLOR_AI_BUBBLE = ("#FFFFFF", "#2B2B2B")
COLOR_BG = ("#F2F2F2", "#1a1a1a")
COLOR_SIDEBAR = ("#EBEBEB", "#212121")

# Linux ç¨³å®šæ€§: å…³é—­è‡ªåŠ¨ DPI ç¼©æ”¾é˜²æ­¢å´©æºƒ
if platform.system() == "Linux":
    ctk.deactivate_automatic_dpi_awareness()

ctk.set_appearance_mode("System")
ctk.set_default_color_theme("blue")

# --- è¾…åŠ©å‡½æ•° ---
def get_base_path():
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))

class AttachmentChip(ctk.CTkFrame):
    def __init__(self, master, filename, command_delete, **kwargs):
        super().__init__(master, fg_color=("gray85", "gray30"), corner_radius=10, **kwargs)
        lbl = ctk.CTkLabel(self, text=filename, font=("Arial", 11))
        lbl.pack(side="left", padx=(10, 5), pady=2)
        btn = ctk.CTkButton(self, text="Ã—", width=20, height=20, 
                            fg_color="transparent", hover_color=("gray70", "gray40"),
                            text_color="red", font=("Arial", 14, "bold"),
                            command=command_delete)
        btn.pack(side="right", padx=(0, 5), pady=2)

class ChatBubble(ctk.CTkFrame):
    def __init__(self, master, role, text="", is_reasoning=False, timestamp=None, is_streaming=False, **kwargs):
        super().__init__(master, fg_color="transparent", **kwargs)
        self.role = role
        self.raw_text = text 
        self.is_reasoning = is_reasoning
        self.is_streaming = is_streaming
        
        self.grid_columnconfigure(0 if role == "user" else 1, weight=1)
        self.grid_columnconfigure(1 if role == "user" else 0, weight=0)
        
        if role == "user":
            bubble_color = COLOR_USER_BUBBLE
            text_color = "black"
            anchor = "e"
        else:
            bubble_color = COLOR_AI_BUBBLE
            text_color = ("black", "white")
            anchor = "w"

        if is_reasoning:
            bubble_color = ("#F0F0F0", "#333333")
            self.text_color_val = "gray"
            self.prefix = "ğŸ§  æ·±åº¦æ€è€ƒ:\n"
        else:
            self.text_color_val = ("black", "white")
            self.prefix = ""

        self.bubble_inner = ctk.CTkFrame(self, fg_color=bubble_color, corner_radius=12)
        self.bubble_inner.grid(row=0, column=1 if role == "user" else 0, padx=10, pady=5, sticky=anchor)

        self.content_frame = ctk.CTkFrame(self.bubble_inner, fg_color="transparent")
        self.content_frame.pack(fill="both", padx=10, pady=10)

        # å­—ä½“ä¿®å¤ï¼šLinuxä¸‹ä½¿ç”¨ Arial ä»¥é˜² Microsoft YaHei UI ç¼ºå¤±å¯¼è‡´æ–‡å­—ä¸å¯è§
        self.main_font = ("Arial", 14) 

        if self.is_streaming:
            self.stream_widget = ctk.CTkTextbox(
                self.content_frame, 
                font=self.main_font, 
                text_color=self.text_color_val,
                fg_color="transparent", 
                wrap="word",
                height=40, 
                width=300
            )
            self.stream_widget.pack(fill="both", expand=True)
            self.stream_widget.insert("0.0", self.prefix + text)
            self.stream_widget.configure(state="disabled")
        else:
            self.render_final_content(self.prefix + text)

        self.bottom_bar = ctk.CTkFrame(self.bubble_inner, fg_color="transparent", height=20)
        self.bottom_bar.pack(fill="x", padx=10, pady=(0, 5))
        
        self.btn_copy = ctk.CTkButton(self.bottom_bar, text="ğŸ“‹ å¤åˆ¶", width=50, height=20,
                                      fg_color="transparent", hover_color=("gray80", "gray40"),
                                      text_color="gray", font=("Arial", 10),
                                      command=self.copy_content)
        self.btn_copy.pack(side="right")

        if timestamp:
            ctk.CTkLabel(self.bottom_bar, text=timestamp, font=("Arial", 10), text_color="gray").pack(side="left")

    def append_stream_text(self, delta_text):
        if not self.is_streaming: return
        self.raw_text += delta_text
        self.stream_widget.configure(state="normal")
        self.stream_widget.insert("end", delta_text)
        self.stream_widget.configure(state="disabled")
        self.stream_widget.see("end")

    def finish_stream(self):
        if not self.is_streaming: return
        self.is_streaming = False
        self.stream_widget.destroy()
        self.render_final_content(self.prefix + self.raw_text)

    def copy_content(self):
        try:
            content = self.raw_text
            if not content: return
            if pyperclip:
                try: pyperclip.copy(content)
                except:
                    self.master.clipboard_clear()
                    self.master.clipboard_append(content)
                    self.master.update()
            else:
                self.master.clipboard_clear()
                self.master.clipboard_append(content)
                self.master.update()

            self.btn_copy.configure(text="âœ… æˆåŠŸ", text_color="green")
            self.after(1500, lambda: self.btn_copy.configure(text="ğŸ“‹ å¤åˆ¶", text_color="gray"))
        except Exception:
            self.btn_copy.configure(text="âŒ å¤±è´¥", text_color="red")

    def render_final_content(self, text):
        parts = re.split(r'(```[\s\S]*?```)', text)
        for part in parts:
            if part.startswith("```") and part.endswith("```"):
                code = part.strip("`")
                if '\n' in code: code = code.split('\n', 1)[1]
                
                f = ctk.CTkFrame(self.content_frame, fg_color="#1E1E1E", corner_radius=5)
                f.pack(fill="x", pady=5)
                
                # ä»£ç å­—ä½“ä½¿ç”¨ Courier æˆ– Consolas
                t = ctk.CTkTextbox(f, font=("Courier", 12), text_color="#D4D4D4", fg_color="transparent", 
                                   height=min(len(code.split('\n'))*20 + 20, 400), wrap="none")
                t.insert("0.0", code)
                t.configure(state="disabled")
                t.pack(fill="x", padx=5, pady=5)
                
                def copy_code(c=code):
                    if pyperclip: pyperclip.copy(c)
                    else: 
                        self.master.clipboard_clear()
                        self.master.clipboard_append(c)
                        self.master.update()
                
                ctk.CTkButton(f, text="å¤åˆ¶ä»£ç ", height=20, width=60, font=("Arial", 10),
                              fg_color="#333333", hover_color="#444444",
                              command=copy_code).pack(anchor="ne", padx=5, pady=2)
            else:
                if part:
                    ctk.CTkLabel(self.content_frame, text=part, text_color=self.text_color_val, justify="left", 
                                 font=self.main_font, wraplength=600).pack(fill="x", anchor="w")

class DeepSeekApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title(f"{APP_NAME} {APP_VERSION}")
        self.geometry("1300x850")
        
        self.base_dir = get_base_path()
        self.config_path = os.path.join(self.base_dir, "config.json")
        self.history_path = os.path.join(self.base_dir, "sessions.json")

        self.config = self.load_json(self.config_path, DEFAULT_CONFIG)
        self.sessions = self.load_json(self.history_path, [])
        
        if not self.sessions or not isinstance(self.sessions, list):
            self.create_new_session(save=False)
        else:
            self.current_session_index = 0
            
        self.attachments = []
        self.client = None
        self.is_running = False
        self.last_scroll_time = 0

        self.setup_ui()
        self.after(200, self.load_current_session_ui)
        self.update_model_status_display()
        
        if self.config["api_key"]:
            self.init_client()

    def load_json(self, path, default):
        if os.path.exists(path):
            try: return json.load(open(path, "r", encoding="utf-8"))
            except: pass
        return default

    def save_config(self):
        try: json.dump(self.config, open(self.config_path, "w", encoding="utf-8"), indent=2)
        except: pass

    def save_sessions(self):
        try: json.dump(self.sessions, open(self.history_path, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
        except: pass

    def init_client(self):
        if not self.config["api_key"]: return
        self.client = OpenAI(api_key=self.config["api_key"], base_url="https://api.deepseek.com")

    # --- UI æ„å»º ---
    def setup_ui(self):
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        # === å·¦ä¾§ ===
        self.sidebar = ctk.CTkFrame(self, width=250, corner_radius=0, fg_color=COLOR_SIDEBAR)
        self.sidebar.grid(row=0, column=0, sticky="nsew")
        self.sidebar.grid_rowconfigure(4, weight=1) 

        # 1. Header
        top_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        top_frame.grid(row=0, column=0, sticky="ew", padx=15, pady=(25, 15))
        ctk.CTkLabel(top_frame, text=APP_NAME, font=("Arial", 22, "bold")).pack(anchor="w")
        dev_frame = ctk.CTkFrame(top_frame, fg_color="transparent")
        dev_frame.pack(anchor="w", pady=(3, 0))
        ctk.CTkLabel(dev_frame, text="Developer:", font=("Arial", 11, "bold"), text_color="gray60").pack(side="left")
        ctk.CTkLabel(dev_frame, text=DEV_NAME, font=("Arial", 11, "bold"), text_color="#3498DB").pack(side="left", padx=5)
        ctk.CTkLabel(top_frame, text=APP_VERSION, font=("Arial", 10), text_color="gray50").pack(anchor="w", pady=(2,0))

        # 2. New Chat
        self.btn_new = ctk.CTkButton(self.sidebar, text="+ å¼€å¯æ–°å¯¹è¯", height=40, font=("Arial", 14), 
                                     fg_color="#3498DB", hover_color="#2980B9",
                                     command=lambda: self.create_new_session(save=True))
        self.btn_new.grid(row=1, column=0, padx=15, pady=(0, 10), sticky="ew")

        # 3. Status
        self.status_frame = ctk.CTkFrame(self.sidebar, fg_color=("white", "#333333"), corner_radius=8)
        self.status_frame.grid(row=2, column=0, sticky="ew", padx=15, pady=5)
        ctk.CTkLabel(self.status_frame, text="å½“å‰æ¨¡å‹çŠ¶æ€", font=("Arial", 10, "bold"), text_color="gray").pack(pady=(5,0))
        self.lbl_model_status = ctk.CTkLabel(self.status_frame, text="åˆå§‹åŒ–ä¸­...", font=("Arial", 12), text_color="#3498DB")
        self.lbl_model_status.pack(pady=(0,5))

        # 4. History
        ctk.CTkLabel(self.sidebar, text="å†å²è®°å½•", font=("Arial", 12), text_color="gray").grid(row=3, column=0, sticky="nw", padx=15, pady=(10,0))
        self.history_list = ctk.CTkScrollableFrame(self.sidebar, fg_color="transparent")
        self.history_list.grid(row=4, column=0, sticky="nsew", padx=5, pady=5)
        self.render_history_list()

        # 5. Settings
        setting_frame = ctk.CTkFrame(self.sidebar, fg_color=("white", "#2B2B2B"), corner_radius=10)
        setting_frame.grid(row=5, column=0, sticky="ew", padx=10, pady=20)
        
        self.r1_var = ctk.BooleanVar(value=self.config["is_r1"])
        ctk.CTkSwitch(setting_frame, text="æ·±åº¦æ€è€ƒ (R1)", variable=self.r1_var, command=self.update_settings).pack(pady=5, padx=10, anchor="w")
        self.search_var = ctk.BooleanVar(value=self.config["use_search"])
        ctk.CTkSwitch(setting_frame, text="è”ç½‘æœç´¢", variable=self.search_var, command=self.update_settings).pack(pady=5, padx=10, anchor="w")

        self.entry_key = ctk.CTkEntry(setting_frame, placeholder_text="API Key", show="*")
        self.entry_key.insert(0, self.config["api_key"])
        self.entry_key.pack(pady=5, padx=10, fill="x")
        ctk.CTkButton(setting_frame, text="ä¿å­˜é…ç½®", height=24, command=self.save_key).pack(pady=10)

        # 6. Clear
        self.btn_clear = ctk.CTkButton(self.sidebar, text="ğŸ—‘ï¸ æ¸…ç©ºæ‰€æœ‰", fg_color="transparent", text_color="#C0392B", hover_color=("#FADBD8", "#522"), command=self.clear_all_history)
        self.btn_clear.grid(row=6, column=0, sticky="ew", padx=15, pady=10)

        # === Right ===
        self.main_area = ctk.CTkFrame(self, fg_color=COLOR_BG)
        self.main_area.grid(row=0, column=1, sticky="nsew")
        self.main_area.grid_rowconfigure(0, weight=1)
        self.main_area.grid_columnconfigure(0, weight=1)

        self.chat_scroll = ctk.CTkScrollableFrame(self.main_area, fg_color="transparent")
        self.chat_scroll.grid(row=0, column=0, sticky="nsew", padx=10, pady=10)

        input_frame = ctk.CTkFrame(self.main_area, fg_color=("white", "#2B2B2B"), height=180)
        input_frame.grid(row=1, column=0, sticky="ew", padx=20, pady=20)
        input_frame.grid_columnconfigure(0, weight=1)

        self.attach_display = ctk.CTkScrollableFrame(input_frame, height=40, orientation="horizontal", fg_color="transparent")
        self.attach_display.grid(row=0, column=0, columnspan=2, sticky="ew", padx=5, pady=5)
        
        self.entry_msg = ctk.CTkTextbox(input_frame, height=80, font=("Arial", 14), fg_color="transparent", border_width=0)
        self.entry_msg.grid(row=1, column=0, sticky="nsew", padx=10, pady=5)
        self.entry_msg.bind("<Return>", self.on_enter_press)

        btn_box = ctk.CTkFrame(input_frame, fg_color="transparent")
        btn_box.grid(row=1, column=1, sticky="s", padx=10, pady=10)
        
        self.btn_attach = ctk.CTkButton(btn_box, text="ğŸ“", width=40, command=self.upload_files)
        self.btn_attach.pack(side="left", padx=2)
        
        self.btn_send = ctk.CTkButton(btn_box, text="å‘é€", width=80, command=self.send_message)
        self.btn_send.pack(side="left", padx=2)
        
        self.btn_stop = ctk.CTkButton(btn_box, text="â¹", width=40, fg_color="#C0392B", command=self.stop_generation)

    # --- Logic ---
    def update_settings(self):
        self.config["use_search"] = self.search_var.get()
        self.config["is_r1"] = self.r1_var.get()
        self.config["model"] = "deepseek-reasoner" if self.r1_var.get() else "deepseek-chat"
        self.save_config()
        self.update_model_status_display()

    def update_model_status_display(self):
        model_name = "DeepSeek-R1 (æ·±åº¦æ¨ç†)" if self.r1_var.get() else "DeepSeek-V3 (æé€Ÿå¯¹è¯)"
        search_status = " + ğŸŒ è”ç½‘" if self.search_var.get() else ""
        self.lbl_model_status.configure(text=f"{model_name}{search_status}")

    def save_key(self):
        key = self.entry_key.get().strip()
        self.config["api_key"] = key
        self.save_config()
        self.init_client()
        messagebox.showinfo("OK", "Key Saved")

    def throttled_scroll_to_bottom(self):
        now = time.time()
        if now - self.last_scroll_time > 0.05:
            self.chat_scroll.update_idletasks()
            try: self.chat_scroll._parent_canvas.yview_moveto(1.0)
            except: pass
            self.last_scroll_time = now

    def force_scroll_to_bottom(self):
        self.chat_scroll.update_idletasks()
        try: self.chat_scroll._parent_canvas.yview_moveto(1.0)
        except: pass

    def send_message(self):
        text = self.entry_msg.get("0.0", "end").strip()
        if not text and not self.attachments: return
        if not self.client: return messagebox.showerror("Error", "No API Key")

        display_text = text
        full_prompt = ""
        
        if self.attachments:
            files_str = "\n".join([f"æ–‡ä»¶[{f['name']}]:\n{f['content']}" for f in self.attachments])
            full_prompt += files_str + "\n\n"
            display_text += f"\n[å·²å‘é€ {len(self.attachments)} ä¸ªæ–‡ä»¶]"
            self.attachments = []
            self.render_attachments_ui()
        
        full_prompt += text
        ts = datetime.now().strftime("%H:%M")

        self.entry_msg.delete("0.0", "end")
        self.add_bubble_ui("user", display_text, timestamp=ts)
        self.force_scroll_to_bottom()

        session = self.sessions[self.current_session_index]
        if len(session["messages"]) == 0:
            session["title"] = text[:15]
            self.render_history_list()
        
        session["messages"].append({"role": "user", "content": full_prompt, "timestamp": ts})
        self.save_sessions()

        self.is_running = True
        self.btn_send.pack_forget()
        self.btn_stop.pack(side="left")
        threading.Thread(target=self.process_stream, args=(full_prompt,), daemon=True).start()

    def process_stream(self, prompt):
        if self.search_var.get():
            self.after(0, lambda: self.add_bubble_ui("ai", "ğŸ” æ­£åœ¨æœç´¢...", timestamp="System"))
            s = self.perform_search(prompt[-100:])
            if s: prompt = f"å‚è€ƒèµ„æ–™:\n{s}\n\né—®é¢˜:\n{prompt}"

        session = self.sessions[self.current_session_index]
        api_msgs = [{"role": "system", "content": self.config["system_prompt"]}]
        for m in session["messages"][-6:]:
            api_msgs.append({"role": "user" if m["role"]=="user" else "assistant", "content": m["content"]})
        if api_msgs[-1]["content"] != prompt:
             api_msgs.append({"role": "user", "content": prompt})

        try:
            response = self.client.chat.completions.create(
                model=self.config["model"],
                messages=api_msgs,
                stream=True
            )
            
            r1_text = ""
            ai_text = ""
            bubble_r1 = None
            bubble_ai = None
            
            def get_r1():
                nonlocal bubble_r1
                if not bubble_r1: bubble_r1 = self.add_bubble_ui("ai", "", is_reasoning=True, is_streaming=True)
                return bubble_r1
            def get_ai():
                nonlocal bubble_ai
                if not bubble_ai: bubble_ai = self.add_bubble_ui("ai", "", is_streaming=True)
                return bubble_ai

            for chunk in response:
                if not self.is_running: break
                delta = chunk.choices[0].delta
                
                if hasattr(delta, 'reasoning_content') and delta.reasoning_content:
                    c = delta.reasoning_content # ä¿®å¤å˜é‡å®šä¹‰
                    r1_text += c
                    self.after(0, lambda b=get_r1(), t=c: b.append_stream_text(t))
                    self.after(0, self.throttled_scroll_to_bottom)

                if hasattr(delta, 'content') and delta.content:
                    c = delta.content # ä¿®å¤æ ¸å¿ƒBug: v2.5.1 æ­¤å¤„æ¼æ‰äº†cçš„å®šä¹‰
                    ai_text += c
                    self.after(0, lambda b=get_ai(), t=c: b.append_stream_text(t))
                    self.after(0, self.throttled_scroll_to_bottom)

            ts = datetime.now().strftime("%H:%M")
            session["messages"].append({"role": "ai", "content": ai_text, "reasoning": r1_text, "timestamp": ts})
            self.save_sessions()

        except Exception as e:
            self.after(0, lambda: messagebox.showerror("API Error", str(e)))
        
        finally:
            self.is_running = False
            self.after(0, self.reset_ui)
            self.after(0, self.force_scroll_to_bottom)
            # ç»“æŸåå°†æµå¼æ–‡æœ¬æ¡†è½¬ä¸ºé™æ€Markdown
            self.after(0, lambda: bubble_ai.finish_stream() if bubble_ai else None)
            self.after(0, lambda: bubble_r1.finish_stream() if bubble_r1 else None)

    # ... (å…¶ä½™æ–¹æ³•ä¿æŒä¸å˜) ...
    def create_new_session(self, save=True):
        new_session = {"id": str(uuid.uuid4()), "title": "æ–°å¯¹è¯", "time": datetime.now().strftime("%m-%d"), "messages": []}
        self.sessions.insert(0, new_session)
        self.current_session_index = 0
        if save:
            self.save_sessions()
            self.render_history_list()
            self.load_current_session_ui()

    def switch_session(self, index):
        self.current_session_index = index
        self.render_history_list()
        self.load_current_session_ui()

    def delete_session(self, index):
        if len(self.sessions) <= 1:
            self.create_new_session(save=False)
            self.sessions = [self.sessions[0]]
        else:
            del self.sessions[index]
            if self.current_session_index >= index: self.current_session_index = max(0, self.current_session_index - 1)
        self.save_sessions()
        self.render_history_list()
        self.load_current_session_ui()

    def render_history_list(self):
        for widget in self.history_list.winfo_children(): widget.destroy()
        for i, session in enumerate(self.sessions):
            color = ("#D1D1D1", "#3A3A3A") if i == self.current_session_index else "transparent"
            item = ctk.CTkFrame(self.history_list, fg_color=color, corner_radius=6)
            item.pack(fill="x", pady=2)
            item.bind("<Button-1>", lambda e, idx=i: self.switch_session(idx))
            
            title = session.get("title", "æ— æ ‡é¢˜")
            if len(title) > 12: title = title[:12] + "..."
            lbl_title = ctk.CTkLabel(item, text=title, font=("Arial", 13, "bold"))
            lbl_title.pack(anchor="w", padx=10, pady=(5,0))
            lbl_title.bind("<Button-1>", lambda e, idx=i: self.switch_session(idx))
            
            btn_del = ctk.CTkButton(item, text="Ã—", width=20, height=20, fg_color="transparent", text_color="gray", hover_color="red", command=lambda idx=i: self.delete_session(idx))
            btn_del.place(relx=1.0, rely=0.5, anchor="e", x=-5)

    def load_current_session_ui(self):
        for widget in self.chat_scroll.winfo_children(): widget.destroy()
        self.attachments = []
        self.render_attachments_ui()
        session = self.sessions[self.current_session_index]
        for msg in session.get("messages", []):
            if msg["role"] == "user": 
                self.add_bubble_ui("user", msg["content"], timestamp=msg.get("timestamp"))
            else:
                if msg.get("reasoning"): 
                    self.add_bubble_ui("ai", msg["reasoning"], is_reasoning=True, timestamp=msg.get("timestamp"), is_streaming=False)
                if msg.get("content"): 
                    self.add_bubble_ui("ai", msg["content"], is_reasoning=False, timestamp=msg.get("timestamp"), is_streaming=False)
        self.force_scroll_to_bottom()

    def upload_files(self):
        files = filedialog.askopenfilenames()
        if not files: return
        for path in files:
            self.attachments.append({"name": os.path.basename(path), "content": self.extract_text(path)})
        self.render_attachments_ui()

    def render_attachments_ui(self):
        for w in self.attach_display.winfo_children(): w.destroy()
        for i, att in enumerate(self.attachments):
            AttachmentChip(self.attach_display, att["name"], lambda idx=i: self.remove_attachment(idx)).pack(side="left", padx=5)

    def remove_attachment(self, index):
        del self.attachments[index]
        self.render_attachments_ui()

    def extract_text(self, filepath):
        try:
            ext = os.path.splitext(filepath)[1].lower()
            if ext == '.pdf':
                reader = pypdf.PdfReader(filepath)
                return "\n".join([p.extract_text() or "" for p in reader.pages])
            elif ext == '.docx':
                doc = Document(filepath)
                return "\n".join([p.text for p in doc.paragraphs])
            elif ext in ['.xlsx', '.xls', '.csv']:
                if pd:
                    df = pd.read_excel(filepath) if 'xls' in ext else pd.read_csv(filepath)
                    return df.to_string()
                else: return "[Error: Pandas not installed]"
            elif ext == '.pptx':
                if Presentation:
                    prs = Presentation(filepath)
                    txt = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"): txt.append(shape.text)
                    return "\n".join(txt)
                else: return "[Error: python-pptx not installed]"
            else:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f: return f.read()[:30000]
        except Exception as e: return f"[Read Error: {e}]"

    def clear_all_history(self):
        if messagebox.askyesno("Confirm", "Delete ALL history?"):
            self.sessions = []
            self.create_new_session()

    def add_bubble_ui(self, role, text, is_reasoning=False, timestamp=None, is_streaming=False):
        b = ChatBubble(self.chat_scroll, role, text, is_reasoning, timestamp, is_streaming)
        b.pack(fill="x", pady=5)
        return b

    def reset_ui(self):
        self.btn_stop.pack_forget()
        self.btn_send.pack(side="left", padx=2)
        self.btn_send.configure(state="normal", text="å‘é€")

    def perform_search(self, query):
        try:
            with DDGS() as ddgs:
                r = list(ddgs.text(query, max_results=3))
                if r: return "\n".join([f"- {x['title']}: {x['body']}" for x in r])
        except: pass
        return ""

    def on_enter_press(self, event):
        if not event.state & 0x0001: 
            self.send_message()
            return "break"

    def stop_generation(self):
        self.is_running = False
        self.reset_ui()

if __name__ == "__main__":
    try:
        app = DeepSeekApp()
        app.mainloop()
    except Exception as e:
        with open("crash_log.txt", "w") as f:
            f.write(traceback.format_exc())
